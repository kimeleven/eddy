from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 컬러 팔레트 ──────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x1B, 0x2A)   # 딥 네이비
ACCENT   = RGBColor(0x00, 0xC2, 0xFF)   # 시안 블루
ACCENT2  = RGBColor(0xFF, 0x6B, 0x35)   # 오렌지 (경고/강조)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GRAY     = RGBColor(0x8A, 0x9B, 0xAE)
GREEN    = RGBColor(0x2E, 0xCC, 0x71)
RED      = RGBColor(0xE7, 0x4C, 0x3C)
YELLOW   = RGBColor(0xF3, 0x9C, 0x12)
CARD_BG  = RGBColor(0x16, 0x2B, 0x40)   # 카드 배경

W = Inches(13.33)   # 와이드 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]   # 완전 빈 레이아웃

# ── 헬퍼 ─────────────────────────────────────────────────────
def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    return shp

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.color.rgb = color
    return tb

def accent_bar(slide, y=Inches(0.55), w_ratio=0.18):
    add_rect(slide, Inches(0.5), y, Inches(13.33 * w_ratio), Inches(0.06), ACCENT)

def slide_number(slide, n):
    add_text(slide, f"{n:02d}", Inches(12.5), Inches(7.1),
             Inches(0.8), Inches(0.3), size=11, color=GRAY, align=PP_ALIGN.RIGHT)

def chip(slide, label, l, t, color=ACCENT):
    add_rect(slide, l, t, Inches(1.5), Inches(0.32), color)
    add_text(slide, label, l, t + Inches(0.03),
             Inches(1.5), Inches(0.3), size=11, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

def status_dot(slide, ok: bool, l, t):
    c = GREEN if ok else RED
    add_rect(slide, l, t, Inches(0.18), Inches(0.18), c)

# ════════════════════════════════════════════════════════════
# SLIDE 1 – 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)

# 왼쪽 사이드 액센트 바
add_rect(s, Inches(0), Inches(0), Inches(0.12), H, ACCENT)

# 대제목
add_text(s, "ELDO", Inches(0.5), Inches(1.6),
         Inches(7), Inches(1.5), size=90, bold=True,
         color=WHITE, align=PP_ALIGN.LEFT)

# 서브타이틀
add_text(s, "글로벌 투자 데이터 플랫폼", Inches(0.55), Inches(3.1),
         Inches(7), Inches(0.5), size=22, color=ACCENT, align=PP_ALIGN.LEFT)

# 설명
add_text(s,
    "현황 분석 및 자동화 로드맵\n현재 상태 → 해야 할 일 → 우선순위",
    Inches(0.55), Inches(3.75),
    Inches(7), Inches(0.9), size=16, color=GRAY, align=PP_ALIGN.LEFT)

# 오른쪽 장식 카드
add_rect(s, Inches(8.8), Inches(1.2), Inches(4.0), Inches(5.2), CARD_BG)
add_rect(s, Inches(8.8), Inches(1.2), Inches(4.0), Inches(0.07), ACCENT)

stats = [
    ("프레임워크 / UI / API", "✅  구축 완료"),
    ("DB 스키마",             "✅  정의 완료"),
    ("실제 데이터",           "❌  전혀 없음"),
    ("데이터 파이프라인",     "❌  미구현"),
    ("분석 파일 생성",        "❌  미구현"),
    ("CI/CD",                 "⚠️  배포만 있음"),
    ("미완성 페이지",         "⚠️  3개 페이지"),
]
for i, (k, v) in enumerate(stats):
    yy = Inches(1.55) + i * Inches(0.67)
    add_text(s, k, Inches(9.0), yy, Inches(2.2), Inches(0.5), size=13, color=GRAY)
    color = GREEN if "✅" in v else (RED if "❌" in v else YELLOW)
    add_text(s, v, Inches(11.2), yy, Inches(1.5), Inches(0.5), size=13, bold=True, color=color)

slide_number(s, 1)

# ════════════════════════════════════════════════════════════
# SLIDE 2 – 현재 상태 한눈에
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "현재 상태 한눈에", Inches(0.5), Inches(0.65),
         Inches(10), Inches(0.55), size=28, bold=True, color=WHITE)
add_text(s, "앱 구조는 완성됐지만 데이터가 없어 아무것도 동작하지 않음",
         Inches(0.5), Inches(1.25), Inches(11), Inches(0.4),
         size=15, color=GRAY)

# 6개 카드 2×3
cards = [
    ("Next.js 16\n+ React 19",       "UI / 라우팅",       "✅ 완료",  GREEN),
    ("API Routes\n(8개 엔드포인트)", "서버 API",          "✅ 완료",  GREEN),
    ("Prisma 7\n+ PostgreSQL 18",     "DB 스키마",         "✅ 완료",  GREEN),
    ("Corps / Indicators\nStatements / StockTrades",
                                      "DB 데이터",         "❌ 비어있음", RED),
    ("data/analysis/\n*.json 파일들", "섹터 분석 차트",    "❌ 없음",  RED),
    ("GitLab CI/CD",                  "빌드·배포 자동화",  "⚠️ 배포만", YELLOW),
]

cols, rows = 3, 2
cw, ch = Inches(3.8), Inches(2.2)
mx, my = Inches(0.5), Inches(1.85)
gap_x, gap_y = Inches(0.27), Inches(0.22)

for i, (title, sub, status, sc) in enumerate(cards):
    col = i % cols
    row = i // cols
    lx = mx + col * (cw + gap_x)
    ty = my + row * (ch + gap_y)

    add_rect(s, lx, ty, cw, ch, CARD_BG)
    add_rect(s, lx, ty, cw, Inches(0.05), sc)          # 상단 컬러 선

    add_text(s, title, lx + Inches(0.18), ty + Inches(0.15),
             cw - Inches(0.3), Inches(0.8), size=14, bold=True, color=WHITE)
    add_text(s, sub, lx + Inches(0.18), ty + Inches(0.95),
             cw - Inches(0.3), Inches(0.35), size=11, color=GRAY)
    add_text(s, status, lx + Inches(0.18), ty + Inches(1.55),
             cw - Inches(0.3), Inches(0.4), size=13, bold=True, color=sc)

slide_number(s, 2)

# ════════════════════════════════════════════════════════════
# SLIDE 3 – 핵심 문제: 데이터 파이프라인
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "핵심 문제  —  데이터 파이프라인 전무", Inches(0.5), Inches(0.65),
         Inches(12), Inches(0.55), size=26, bold=True, color=WHITE)
add_text(s, "DB 테이블은 모두 정의되어 있으나 데이터를 채우는 스크립트가 하나도 없음",
         Inches(0.5), Inches(1.25), Inches(12), Inches(0.4),
         size=14, color=GRAY)

# 흐름 화살표 다이어그램
boxes = [
    ("외부 데이터 소스\n(KRX / DART / Yahoo\nFinance 등)", ACCENT2),
    ("수집 스크립트\n(미구현)", RED),
    ("PostgreSQL DB\n(스키마만 있음)", RED),
    ("분석 JSON 생성\n(미구현)", RED),
    ("웹 앱 차트\n(섹터 페이지)", ACCENT),
]

bw, bh = Inches(2.0), Inches(1.1)
by = Inches(2.2)
start_x = Inches(0.45)
gap = Inches(0.55)

for i, (label, color) in enumerate(boxes):
    lx = start_x + i * (bw + gap)
    add_rect(s, lx, by, bw, bh, CARD_BG)
    add_rect(s, lx, by, bw, Inches(0.05), color)
    add_text(s, label, lx + Inches(0.1), by + Inches(0.12),
             bw - Inches(0.2), bh - Inches(0.15),
             size=12, bold=False, color=WHITE if color != RED else RGBColor(0xFF,0xAA,0xAA))

    if i < len(boxes) - 1:
        ax = lx + bw + Inches(0.12)
        ay = by + bh / 2 - Inches(0.05)
        add_rect(s, ax, ay, Inches(0.38), Inches(0.06), GRAY)
        add_text(s, "▶", ax + Inches(0.22), ay - Inches(0.18),
                 Inches(0.3), Inches(0.4), size=14, color=GRAY)

# 빠진 것들 목록
add_text(s, "빠진 것들", Inches(0.5), Inches(3.65),
         Inches(5), Inches(0.4), size=16, bold=True, color=ACCENT)

missing = [
    "기업 마스터 수집  (corps, corps_emsec 테이블)",
    "재무제표 수집  (statements / us_statements)",
    "재무지표 계산 및 적재  (P/E, EV/EBITDA 등 100+ 지표)",
    "일별 주가 수집  (stock_trades — OHLCV, 시가총액)",
    "EMSEC 업종 계층 초기 적재  (섹터 트리)",
    "섹터 분석 JSON 생성 배치  (차트 데이터 파일)",
    "스케줄링  (매일 자동 수집)",
]

for i, txt in enumerate(missing):
    yy = Inches(4.1) + i * Inches(0.42)
    add_rect(s, Inches(0.5), yy + Inches(0.08), Inches(0.12), Inches(0.12), ACCENT2)
    add_text(s, txt, Inches(0.75), yy, Inches(9), Inches(0.38), size=13, color=WHITE)

slide_number(s, 3)

# ════════════════════════════════════════════════════════════
# SLIDE 4 – DB·인프라·CI/CD 자동화
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "DB · 인프라 · CI/CD 자동화", Inches(0.5), Inches(0.65),
         Inches(12), Inches(0.55), size=26, bold=True, color=WHITE)

# 왼쪽 컬럼 – DB / 인프라
add_rect(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(5.7), CARD_BG)
add_rect(s, Inches(0.5), Inches(1.35), Inches(5.8), Inches(0.05), ACCENT)
add_text(s, "DB & 인프라", Inches(0.7), Inches(1.45),
         Inches(5.4), Inches(0.45), size=16, bold=True, color=ACCENT)

db_items = [
    ("prisma/seed.ts 작성",          "EMSEC, 기준 코드 초기 적재",       "❌"),
    ("scripts/convert-schema.ts",    "package.json에 등록됐지만 파일 없음", "❌"),
    (".env.example 작성",             "현재 환경변수 문서화 전무",          "❌"),
    ("DB 셋업 원커맨드화",            "compose up → migrate → seed 한 번에","❌"),
    ("앱 컨테이너 추가",              "compose.yml에 Next.js 서비스 없음",  "❌"),
    ("로컬 샘플 데이터 seed",         "새 개발자 온보딩용 더미 데이터",     "❌"),
]

for i, (title, desc, st) in enumerate(db_items):
    yy = Inches(2.0) + i * Inches(0.72)
    add_text(s, st, Inches(0.65), yy, Inches(0.3), Inches(0.5), size=13, color=RED)
    add_text(s, title, Inches(1.0), yy, Inches(5.0), Inches(0.3),
             size=13, bold=True, color=WHITE)
    add_text(s, desc, Inches(1.0), yy + Inches(0.28), Inches(5.0), Inches(0.3),
             size=11, color=GRAY)

# 오른쪽 컬럼 – CI/CD
add_rect(s, Inches(6.7), Inches(1.35), Inches(6.1), Inches(5.7), CARD_BG)
add_rect(s, Inches(6.7), Inches(1.35), Inches(6.1), Inches(0.05), YELLOW)
add_text(s, "CI/CD 파이프라인", Inches(6.9), Inches(1.45),
         Inches(5.7), Inches(0.45), size=16, bold=True, color=YELLOW)

# 현재 vs 목표
add_text(s, "현재", Inches(6.9), Inches(2.0),
         Inches(2), Inches(0.35), size=13, bold=True, color=RED)
add_rect(s, Inches(6.9), Inches(2.4), Inches(2.5), Inches(0.5), RGBColor(0x2A,0x10,0x10))
add_text(s, "deploy (webhook만)", Inches(7.0), Inches(2.45),
         Inches(2.3), Inches(0.4), size=12, color=RGBColor(0xFF,0x88,0x88))

add_text(s, "목표", Inches(9.6), Inches(2.0),
         Inches(2), Inches(0.35), size=13, bold=True, color=GREEN)

stages = [
    ("lint",     ACCENT,  Inches(9.6), Inches(2.4)),
    ("build",    ACCENT,  Inches(9.6), Inches(3.05)),
    ("migrate",  ACCENT,  Inches(9.6), Inches(3.7)),
    ("test",     ACCENT,  Inches(9.6), Inches(4.35)),
    ("deploy",   GREEN,   Inches(9.6), Inches(5.0)),
]
for label, col, lx, ty in stages:
    add_rect(s, lx, ty, Inches(2.6), Inches(0.45), CARD_BG)
    add_rect(s, lx, ty, Inches(0.06), Inches(0.45), col)
    add_text(s, label, lx + Inches(0.15), ty + Inches(0.08),
             Inches(2.3), Inches(0.35), size=13, bold=True, color=WHITE)

slide_number(s, 4)

# ════════════════════════════════════════════════════════════
# SLIDE 5 – 미완성 페이지 & 우선순위 로드맵
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
bg(s)
accent_bar(s)

add_text(s, "미완성 페이지 & 우선순위 로드맵", Inches(0.5), Inches(0.65),
         Inches(12), Inches(0.55), size=26, bold=True, color=WHITE)

# 미완성 페이지 3개
pages = [
    ("/valuation", "하드코딩 더미 데이터\n→ DB 연동 + 밸류에이션 로직 필요", ACCENT2),
    ("/peer",      "텍스트 placeholder\n→ 동종업체 탐색 UI 전체 구현 필요",  YELLOW),
    ("/compare",   "텍스트 placeholder\n→ 멀티기업 비교 대시보드 구현 필요", YELLOW),
]

for i, (route, desc, col) in enumerate(pages):
    lx = Inches(0.5) + i * Inches(4.18)
    add_rect(s, lx, Inches(1.35), Inches(3.9), Inches(1.6), CARD_BG)
    add_rect(s, lx, Inches(1.35), Inches(3.9), Inches(0.05), col)
    add_text(s, route, lx + Inches(0.18), Inches(1.45),
             Inches(3.6), Inches(0.4), size=16, bold=True, color=col)
    add_text(s, desc, lx + Inches(0.18), Inches(1.9),
             Inches(3.6), Inches(0.85), size=12, color=GRAY)

# 로드맵 타임라인
add_text(s, "우선순위 로드맵", Inches(0.5), Inches(3.2),
         Inches(10), Inches(0.45), size=16, bold=True, color=ACCENT)

roadmap = [
    ("1", "데이터 소스 확정 및 기업·주가 수집 스크립트 작성",           "이게 없으면 아무것도 동작 안 함",    ACCENT2),
    ("2", "EMSEC + 기준 데이터 시드 (업종 계층, 기준 코드)",            "company / sectors 페이지 기본 동작", ACCENT2),
    ("3", "재무제표 + 지표 수집 파이프라인",                             "기업 상세 페이지, 섹터 차트 데이터", YELLOW),
    ("4", "분석 JSON 생성 배치 + .env.example + DB 자동화",              "재현 가능한 개발 환경 구성",         YELLOW),
    ("5", "CI/CD 보강 + Docker Compose 완성 + 미완성 페이지 구현",       "운영 안정성 및 기능 완성",           GREEN),
]

for i, (num, title, reason, col) in enumerate(roadmap):
    yy = Inches(3.75) + i * Inches(0.67)
    # 번호 배지
    add_rect(s, Inches(0.5), yy + Inches(0.05), Inches(0.35), Inches(0.35), col)
    add_text(s, num, Inches(0.5), yy + Inches(0.03),
             Inches(0.35), Inches(0.38), size=13, bold=True,
             color=BG, align=PP_ALIGN.CENTER)
    add_text(s, title, Inches(0.98), yy,
             Inches(7.5), Inches(0.35), size=13, bold=True, color=WHITE)
    add_text(s, reason, Inches(0.98), yy + Inches(0.33),
             Inches(7.5), Inches(0.3), size=11, color=GRAY)

    # 라벨
    phase = ["즉시", "즉시", "단기", "단기", "중기"][i]
    pc = [ACCENT2, ACCENT2, YELLOW, YELLOW, GREEN][i]
    add_rect(s, Inches(8.7), yy + Inches(0.07), Inches(1.0), Inches(0.3), pc)
    add_text(s, phase, Inches(8.7), yy + Inches(0.06),
             Inches(1.0), Inches(0.3), size=11, bold=True,
             color=BG, align=PP_ALIGN.CENTER)

slide_number(s, 5)

# ── 저장 ─────────────────────────────────────────────────────
out = r"C:\Work\ELDO\ELDO_TODO.pptx"
prs.save(out)
print(f"저장 완료: {out}")
